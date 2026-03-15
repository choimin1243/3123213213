import sys
import os
from PyQt5.QtWidgets import *
from PyQt5.QtGui import *
from PyQt5.QtCore import *

from pptx import Presentation
from pptx.util import Inches


class DrawLabel(QLabel):

    def __init__(self):
        super().__init__()
        self.start = None
        self.end = None
        self.rect = None

    def mousePressEvent(self, e):
        self.start = e.pos()
        self.end = e.pos()
        self.update()

    def mouseMoveEvent(self, e):
        self.end = e.pos()
        self.update()

    def mouseReleaseEvent(self, e):
        self.end = e.pos()
        self.rect = QRect(self.start, self.end)
        self.update()

    def paintEvent(self, e):
        super().paintEvent(e)

        if self.start and self.end:
            painter = QPainter(self)
            painter.setPen(QPen(Qt.red, 3))
            painter.drawRect(QRect(self.start, self.end))


class App(QWidget):

    def __init__(self):
        super().__init__()

        self.prs = None
        self.ppt_path = None
        self.slide_width = None
        self.slide_height = None

        self.initUI()

    def initUI(self):

        layout = QVBoxLayout()

        self.btn1 = QPushButton("PPTX 불러오기")
        self.btn1.clicked.connect(self.load_ppt)

        self.btn2 = QPushButton("사진 넣기")
        self.btn2.clicked.connect(self.insert_images)

        self.label = DrawLabel()
        self.label.setFixedSize(960,540)
        self.label.setStyleSheet("background:white")

        layout.addWidget(self.btn1)
        layout.addWidget(self.btn2)
        layout.addWidget(self.label)

        self.setLayout(layout)

        self.setWindowTitle("PPT 이미지 자동 삽입")
        self.resize(1000,650)

    def load_ppt(self):

        file, _ = QFileDialog.getOpenFileName(self,"PPT 선택","","PPTX (*.pptx)")

        if not file:
            return

        self.ppt_path = file
        self.prs = Presentation(file)

        self.slide_width = self.prs.slide_width
        self.slide_height = self.prs.slide_height

        QMessageBox.information(self,"완료","첫 슬라이드 위치에 빨간 박스를 드래그하세요")

    def insert_images(self):

        if not self.prs:
            QMessageBox.warning(self,"경고","먼저 PPT를 불러오세요")
            return

        if not self.label.rect:
            QMessageBox.warning(self,"경고","영역을 먼저 지정하세요")
            return

        files,_ = QFileDialog.getOpenFileNames(
            self,
            "이미지 선택",
            "",
            "Images (*.png *.jpg *.jpeg)"
        )

        if not files:
            return

        rect = self.label.rect

        # label 좌표 → ppt 좌표 변환
        x_ratio = rect.x() / self.label.width()
        y_ratio = rect.y() / self.label.height()

        w_ratio = rect.width() / self.label.width()
        h_ratio = rect.height() / self.label.height()

        left = self.slide_width * x_ratio
        top = self.slide_height * y_ratio

        width = self.slide_width * w_ratio
        height = self.slide_height * h_ratio

        base_slide = self.prs.slides[0]
        layout = base_slide.slide_layout

        for img in files:

            new_slide = self.prs.slides.add_slide(layout)

            new_slide.shapes.add_picture(
                img,
                left,
                top,
                width=width,
                height=height
            )

        # exe 옆에 저장
        exe_dir = os.path.dirname(sys.executable if getattr(sys,'frozen',False) else __file__)

        save_path = os.path.join(exe_dir,"result.pptx")

        self.prs.save(save_path)

        QMessageBox.information(self,"완료",f"저장됨\n{save_path}")


if __name__ == "__main__":

    app = QApplication(sys.argv)
    window = App()
    window.show()
    sys.exit(app.exec_())
